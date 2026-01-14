"""
Document loaders for various file formats.
Supports: PDF, Word, Excel, PowerPoint, Text, Markdown, CSV, JSON, HTML
"""
import os
import re
import csv
import json
from pathlib import Path
from typing import Tuple

import markdown
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup


def load_pdf(file_path: str) -> Tuple[str, dict]:
    """Extract text from PDF file."""
    reader = PdfReader(file_path)
    text_parts = []

    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text_parts.append(page_text)

    return "\n\n".join(text_parts), {
        "file_type": "pdf",
        "page_count": len(reader.pages),
        "filename": os.path.basename(file_path)
    }


def load_docx(file_path: str) -> Tuple[str, dict]:
    """Extract text from Word document (.docx)."""
    doc = DocxDocument(file_path)
    text_parts = []

    # Extract paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                text_parts.append(" | ".join(row_text))

    return "\n\n".join(text_parts), {
        "file_type": "docx",
        "paragraph_count": len(doc.paragraphs),
        "filename": os.path.basename(file_path)
    }


def load_xlsx(file_path: str) -> Tuple[str, dict]:
    """Extract text from Excel file (.xlsx)."""
    wb = load_workbook(file_path, data_only=True)
    text_parts = []
    total_rows = 0

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        text_parts.append(f"[Sheet: {sheet_name}]")

        for row in sheet.iter_rows():
            row_values = []
            for cell in row:
                if cell.value is not None:
                    row_values.append(str(cell.value))
            if row_values:
                text_parts.append(" | ".join(row_values))
                total_rows += 1

    return "\n".join(text_parts), {
        "file_type": "xlsx",
        "sheet_count": len(wb.sheetnames),
        "row_count": total_rows,
        "filename": os.path.basename(file_path)
    }


def load_pptx(file_path: str) -> Tuple[str, dict]:
    """Extract text from PowerPoint file (.pptx)."""
    prs = Presentation(file_path)
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = [f"[Slide {slide_num}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if len(slide_text) > 1:
            text_parts.append("\n".join(slide_text))

    return "\n\n".join(text_parts), {
        "file_type": "pptx",
        "slide_count": len(prs.slides),
        "filename": os.path.basename(file_path)
    }


def load_markdown(file_path: str) -> Tuple[str, dict]:
    """Load and convert markdown file to plain text."""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    html = markdown.markdown(content)
    text = re.sub(r'<[^>]+>', '', html)

    return text, {
        "file_type": "markdown",
        "filename": os.path.basename(file_path)
    }


def load_text(file_path: str) -> Tuple[str, dict]:
    """Load plain text file."""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    return content, {
        "file_type": "text",
        "filename": os.path.basename(file_path)
    }


def load_csv(file_path: str) -> Tuple[str, dict]:
    """Extract text from CSV file."""
    text_parts = []
    row_count = 0

    with open(file_path, 'r', encoding='utf-8') as f:
        reader = csv.reader(f)
        for row in reader:
            if any(cell.strip() for cell in row):
                text_parts.append(" | ".join(cell.strip() for cell in row))
                row_count += 1

    return "\n".join(text_parts), {
        "file_type": "csv",
        "row_count": row_count,
        "filename": os.path.basename(file_path)
    }


def load_json(file_path: str) -> Tuple[str, dict]:
    """Extract text from JSON file."""
    with open(file_path, 'r', encoding='utf-8') as f:
        data = json.load(f)

    def extract_text(obj, prefix=""):
        """Recursively extract text from JSON structure."""
        texts = []
        if isinstance(obj, dict):
            for key, value in obj.items():
                new_prefix = f"{prefix}.{key}" if prefix else key
                texts.extend(extract_text(value, new_prefix))
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                texts.extend(extract_text(item, f"{prefix}[{i}]"))
        elif obj is not None:
            texts.append(f"{prefix}: {obj}")
        return texts

    text_lines = extract_text(data)

    return "\n".join(text_lines), {
        "file_type": "json",
        "filename": os.path.basename(file_path)
    }


def load_html(file_path: str) -> Tuple[str, dict]:
    """Extract text from HTML file."""
    with open(file_path, 'r', encoding='utf-8') as f:
        content = f.read()

    soup = BeautifulSoup(content, 'lxml')

    # Remove script and style elements
    for script in soup(["script", "style"]):
        script.decompose()

    text = soup.get_text(separator='\n')
    # Clean up whitespace
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    text = "\n".join(lines)

    return text, {
        "file_type": "html",
        "filename": os.path.basename(file_path)
    }


def load_document(file_path: str) -> Tuple[str, dict]:
    """
    Load document based on file extension.

    Supported formats:
    - PDF (.pdf)
    - Word (.docx)
    - Excel (.xlsx)
    - PowerPoint (.pptx)
    - Text (.txt, .text)
    - Markdown (.md, .markdown)
    - CSV (.csv)
    - JSON (.json)
    - HTML (.html, .htm)

    Returns:
        Tuple of (text content, metadata dict)

    Raises:
        ValueError: If file type is not supported
    """
    path = Path(file_path)
    extension = path.suffix.lower()

    loaders = {
        '.pdf': load_pdf,
        '.docx': load_docx,
        '.xlsx': load_xlsx,
        '.pptx': load_pptx,
        '.md': load_markdown,
        '.markdown': load_markdown,
        '.txt': load_text,
        '.text': load_text,
        '.csv': load_csv,
        '.json': load_json,
        '.html': load_html,
        '.htm': load_html,
    }

    if extension not in loaders:
        supported = list(loaders.keys())
        raise ValueError(f"Unsupported file type: {extension}. Supported formats: {supported}")

    return loaders[extension](file_path)


def get_supported_extensions() -> list:
    """Return list of supported file extensions."""
    return ['.pdf', '.docx', '.xlsx', '.pptx', '.md', '.markdown', '.txt', '.text', '.csv', '.json', '.html', '.htm']
